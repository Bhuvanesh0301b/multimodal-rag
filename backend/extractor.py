import fitz
import base64
from pathlib import Path
from typing import List, Dict
from groq import Groq

client = Groq()

def extract_content(file_path: str, original_name: str) -> List[Dict]:
    suffix = Path(file_path).suffix.lower()
    if suffix == ".pdf":
        return extract_pdf(file_path, original_name)
    elif suffix in [".png", ".jpg", ".jpeg"]:
        return extract_image(file_path, original_name)
    elif suffix == ".docx":
        return extract_docx(file_path, original_name)
    elif suffix == ".pptx":
        return extract_pptx(file_path, original_name)
    return []

def extract_pdf(file_path: str, source: str) -> List[Dict]:
    chunks = []
    doc = fitz.open(file_path)
    for page_num, page in enumerate(doc, 1):
        text = page.get_text("text").strip()
        if text and len(text) > 50:
            chunks.append({"type": "text", "content": text, "source": source, "page": page_num})
        for table in page.find_tables():
            data = table.extract()
            if data:
                chunks.append({"type": "table", "content": format_table(data), "source": source, "page": page_num})
        for img_ref in page.get_images(full=True):
            try:
                xref = img_ref[0]
                base_image = doc.extract_image(xref)
                img_bytes = base_image["image"]
                if len(img_bytes) < 5000:
                    continue
                ext = base_image["ext"]
                media_type = "image/jpeg" if ext == "jpg" else f"image/{ext}"
                img_b64 = base64.standard_b64encode(img_bytes).decode()
                description = describe_image(img_b64, media_type)
                chunks.append({"type": "image", "content": description, "source": source, "page": page_num})
            except Exception as e:
                print(f"Skipping image on page {page_num}: {e}")
    doc.close()
    return chunks

def extract_image(file_path: str, source: str) -> List[Dict]:
    with open(file_path, "rb") as f:
        img_bytes = f.read()
    ext = Path(file_path).suffix.lower().lstrip(".")
    media_type = "image/jpeg" if ext in ["jpg", "jpeg"] else f"image/{ext}"
    img_b64 = base64.standard_b64encode(img_bytes).decode()
    description = describe_image(img_b64, media_type)
    return [{"type": "image", "content": description, "source": source, "page": 1}]

def extract_docx(file_path: str, source: str) -> List[Dict]:
    try:
        from docx import Document
    except ImportError:
        return []
    chunks = []
    doc = Document(file_path)
    buffer = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            buffer.append(text)
        if len(" ".join(buffer).split()) > 400:
            chunks.append({"type": "text", "content": " ".join(buffer), "source": source, "page": 1})
            buffer = []
    if buffer:
        chunks.append({"type": "text", "content": " ".join(buffer), "source": source, "page": 1})
    for table in doc.tables:
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        chunks.append({"type": "table", "content": format_table(rows), "source": source, "page": 1})
    return chunks

def extract_pptx(file_path: str, source: str) -> List[Dict]:
    try:
        from pptx import Presentation
    except ImportError:
        return []
    chunks = []
    prs = Presentation(file_path)
    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
            if shape.shape_type == 13:
                try:
                    img_bytes = shape.image.blob
                    if len(img_bytes) > 5000:
                        media_type = shape.image.content_type
                        img_b64 = base64.standard_b64encode(img_bytes).decode()
                        description = describe_image(img_b64, media_type)
                        chunks.append({"type": "image", "content": description, "source": source, "page": slide_num})
                except Exception as e:
                    print(f"Skipping image slide {slide_num}: {e}")
        if texts:
            chunks.append({"type": "text", "content": "\n".join(texts), "source": source, "page": slide_num})
    return chunks

def describe_image(img_b64: str, media_type: str) -> str:
    try:
        response = client.chat.completions.create(
            model="meta-llama/llama-4-scout-17b-16e-instruct",
            max_tokens=400,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "image_url", "image_url": {"url": f"data:{media_type};base64,{img_b64}"}},
                    {"type": "text", "text": "Describe this image in detail for a search index. If it's a chart mention chart type, axes, key values, trends. If it's a table list column names and key data. If it's a diagram explain components and relationships. Be specific about numbers and labels you can see."}
                ]
            }]
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Image present (description failed: {e})"

def format_table(rows):
    if not rows:
        return ""
    lines = []
    for i, row in enumerate(rows):
        line = " | ".join(str(c) for c in row)
        lines.append(line)
        if i == 0:
            lines.append("-" * len(line))
    return "\n".join(lines)